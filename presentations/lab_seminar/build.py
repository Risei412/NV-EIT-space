"""Build the NV-EIT lab seminar deck on the supplied 4:3 corporate template.

Design rules the deck follows (from the requester):
  * IMRAD order
  * every content slide carries a HEAD message (one interpretive sentence,
    subject + verb) and a BOTTOM message (one cold factual sentence); reading
    only those two lines on every slide must give the whole story
  * body text never below 16 pt
  * navy + white, red only for emphasis

Four conventions are taken from the requester's own decks:

  * a TITLE line under the head message, prefixed with the section and a
    number ("Result 2: ..."), so the audience always knows where it is;
  * colour spent sparingly -- structure in navy, one meaning per card tint,
    red at most once a slide and only on the number that is the point;
  * figure text at 18-20 pt, which is why the figures come from
    slide_figures.py rather than from the manuscript's figure set;
  * citations on the same screen as the claim, under a hairline at the foot.

The five bands below are fixed for every content slide, so the eye lands in
the same place each time.
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
FIG = HERE / "figures"
OUT = HERE / "nv_eit_seminar.pptx"

# --- colour.  Structure is navy; a card's tint is its meaning; red is the
# emphasis of last resort and appears at most once per slide.
NAVY = RGBColor(0x1C, 0x30, 0x77)      # template lt2 -- the centre colour
BLUE = RGBColor(0x7F, 0x96, 0xC2)      # template dk2
TINT = RGBColor(0xEE, 0xF4, 0xFF)      # neutral card
SOFT = RGBColor(0xF6, 0xF8, 0xFC)      # secondary neutral card
WARN = RGBColor(0xFF, 0xF6, 0xE8)      # a caveat
BAD = RGBColor(0xFF, 0xF1, 0xF1)       # a failure or a wrong turn
RED = RGBColor(0xC0, 0x00, 0x00)       # emphasis only
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Arial"

# 'コンテンツ1-3' -- the layout the requester's own decks use.  Its logo sits
# at x = 8.36, further right than 'コンテンツ1-1', which buys the head
# message another 0.8 in of line.
L_TITLE, L_SECTION, L_CONTENT = 0, 1, 15

SW, SH = 10.0, 7.5
MARGIN = 0.14                  # near-flush, so the figures get the width
HEAD_Y, HEAD_H = 0.11, 0.76    # interpretive message
HEAD_W = 8.18                  # stops clear of the template logo at x = 8.36
TITLE_Y, TITLE_H = 0.94, 0.50  # "Result 2: ..." -- the location cue
BODY_TOP, BODY_BOT = 1.62, 5.28
BOT_Y, BOT_H = 5.38, 0.80      # factual message
RULE_Y = 6.42
REF_Y, REF_H = 6.47, 0.60      # citations; slide number sits at y = 7.15
BODY_W = SW - 2 * MARGIN

HEAD_PT = 20
TITLE_PT = 22
BOT_PT = 20
REF_PT = 10


# --------------------------------------------------------------- helpers
def strip_placeholders(slide, keep_idx=()):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx not in keep_idx:
            ph._element.getparent().remove(ph._element)


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tb, tf


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def run(p, text, *, size=16, bold=False, color=INK, italic=False, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return r


def rich(p, parts, *, size=16):
    """parts: list of (text, {bold, color, italic})"""
    for text, opt in parts:
        run(p, text, size=opt.get("size", size), bold=opt.get("bold", False),
            color=opt.get("color", INK), italic=opt.get("italic", False))


def head(slide, text, parts=None):
    """Interpretive one-sentence message across the top."""
    tb, tf = textbox(slide, MARGIN, HEAD_Y, HEAD_W, HEAD_H,
                     anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if parts:
        rich(p, parts, size=HEAD_PT)
    else:
        run(p, text, size=HEAD_PT, bold=True, color=NAVY)
    return tb


def title(slide, text):
    """Section and number: the audience's position in the talk."""
    tb, tf = textbox(slide, MARGIN, TITLE_Y, SW - 2 * MARGIN, TITLE_H,
                     anchor=MSO_ANCHOR.MIDDLE)
    tb.name = "TitleLine"
    p = tf.paragraphs[0]
    run(p, text, size=TITLE_PT, bold=True, color=INK)
    return tb


def bottom(slide, text, parts=None):
    """Cold factual one-sentence message along the base."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN),
                                  Inches(BOT_Y), Inches(SW - 2 * MARGIN),
                                  Inches(BOT_H))
    card.fill.solid()
    card.fill.fore_color.rgb = TINT
    card.line.fill.background()
    card.shadow.inherit = False
    try:
        card.adjustments[0] = 0.12
    except Exception:
        pass
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if parts:
        rich(p, parts, size=BOT_PT)
    else:
        run(p, text, size=BOT_PT, color=INK)
    return card


def refs(slide, items):
    """Sources for this slide, on this slide, under a hairline."""
    cn = slide.shapes.add_connector(1, Inches(MARGIN), Inches(RULE_Y),
                                    Inches(SW - MARGIN), Inches(RULE_Y))
    cn.line.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    cn.line.width = Pt(0.75)
    tb, tf = textbox(slide, MARGIN, REF_Y, 9.10, REF_H)
    tb.name = "RefBand"
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0))
        run(p, it, size=REF_PT, italic=True, color=GREY)
    return tb


def picture(slide, name, *, top=None, max_h=None, max_w=None, cx=SW / 2):
    path = FIG / name
    im = Image.open(path)
    ar = im.width / im.height
    mw = max_w if max_w else BODY_W
    mh = max_h if max_h else (BODY_BOT - (top if top else BODY_TOP))
    w = mw
    h = w / ar
    if h > mh:
        h = mh
        w = h * ar
    x = cx - w / 2
    y = top if top else (BODY_TOP + ((BODY_BOT - BODY_TOP) - h) / 2)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                                    Inches(w), Inches(h))


def card(slide, x, y, w, h, *, fill=TINT, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),
                                Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.06
    except Exception:
        pass
    sh.text_frame.word_wrap = True
    sh.text_frame.margin_left = sh.text_frame.margin_right = Inches(0.14)
    sh.text_frame.margin_top = sh.text_frame.margin_bottom = Inches(0.09)
    return sh


def bullets(slide, x, y, w, h, items, *, size=16, space=8):
    tb, tf = textbox(slide, x, y, w, h)
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(space)
        if isinstance(it, str):
            run(p, "•  " + it, size=size)
        else:
            run(p, "•  ", size=size)
            rich(p, it, size=size)
    return tb


# --------------------------------------------------------------- deck
prs = Presentation(str(HERE / "template.pptx"))
sld_lst = prs.slides._sldIdLst
for sid in list(sld_lst):
    rId = sid.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
    prs.part.drop_rel(rId)
    sld_lst.remove(sid)

master = prs.slide_masters[0]


def add(layout_idx, keep=()):
    s = prs.slides.add_slide(master.slide_layouts[layout_idx])
    strip_placeholders(s, keep_idx=keep)
    return s


def section(name, n):
    s = add(L_SECTION, keep=(10,))
    tb, tf = textbox(s, 0.73, 3.35, 8.66, 1.1, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, f"{n}.  ", size=28, bold=True, color=BLUE)
    run(p, name, size=28, bold=True, color=NAVY)
    return s


# ---------------------------------------------------------------- 1 title
s = add(L_TITLE)
tb, tf = textbox(s, 0.71, 1.55, 8.81, 1.6, anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]
run(p, "Where NV electromagnetically induced transparency lives,",
    size=28, bold=True, color=NAVY)
p2 = para(tf)
run(p2, "and where it inverts", size=28, bold=True, color=NAVY)
tb, tf = textbox(s, 0.75, 3.60, 6.03, 0.4)
run(tf.paragraphs[0], "Risei Abe", size=18, bold=True, color=INK)
tb, tf = textbox(s, 0.75, 4.05, 6.5, 1.0)
p = tf.paragraphs[0]
run(p, "Department of Electrical and Electronic Engineering", size=16,
    color=GREY)
p = para(tf)
run(p, "Institute of Science Tokyo", size=16, color=GREY)
tb, tf = textbox(s, 5.69, 7.02, 4.15, 0.32)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.RIGHT
run(p, "Group seminar / progress report", size=16, color=GREY)

# ---------------------------------------------------------------- 2 outline
s = add(L_CONTENT)
title(s, "Outline")
items = [
    ("Introduction", "What EIT is, why NV, and what we ask"),
    ("Method", "The no-go criterion and how each operating point is judged"),
    ("Results", "Sanity check, the island, room temperature, the field"),
    ("Discussion", "What to measure, and why room temperature fails"),
    ("Summary", "Operating window fixed; experiment this year"),
]
y = 1.66
for i, (t, d) in enumerate(items, 1):
    c = card(s, MARGIN, y, BODY_W, 0.90)
    tf = c.text_frame
    p = tf.paragraphs[0]
    run(p, f"{i}.  ", size=18, bold=True, color=NAVY)
    run(p, t, size=18, bold=True, color=NAVY)
    p2 = para(tf)
    run(p2, d, size=16, color=INK)
    y += 1.00

# ================================================== INTRODUCTION
section("Introduction", "I")

# ---- I-1 what EIT is
s = add(L_CONTENT)
head(s, "A control laser makes two excitation paths interfere, turning an "
        "absorber transparent.")
title(s, "Introduction 1: What EIT is")
picture(s, "sf_lambda.png", top=BODY_TOP, max_h=3.66, max_w=5.24, cx=2.76)
bullets(s, 5.58, 1.72, 4.28, 3.4, [
    [("Two ground states, one excited state: the ", {}),
     ("Λ", {"bold": True}), (" system.", {})],
    [("The two absorption paths cancel at two-photon resonance.", {})],
    [("Population is trapped in a ", {}),
     ("dark state", {"bold": True, "color": NAVY}),
     (" that the light cannot excite.", {})],
    [("Requires a long-lived ground coherence: it is destroyed by "
      "dephasing.", {})],
])
bottom(s, "In the ideal Λ system probe absorption vanishes at two-photon "
          "resonance when the ground coherence is long-lived.")
refs(s, ["[1] M. Fleischhauer, A. Imamoğlu and J. P. Marangos, "
         "Rev. Mod. Phys. 77, 633 (2005)."])

# ---- I-2 why NV
s = add(L_CONTENT)
head(s, "We turned to NV because the group-IV centres that motivated us have "
        "no settled model.")
title(s, "Introduction 2: Why NV is the test bed")
c1 = card(s, MARGIN, BODY_TOP, 4.72, 2.85, fill=SOFT)
tf = c1.text_frame
run(tf.paragraphs[0], "SiV / SnV — where we started", size=18, bold=True,
    color=GREY)
for t in ["Inversion symmetric: optical line robust against electric noise",
          "Orbital Λ channel, no spin-overlap obstruction",
          "Room-temperature EIT was the hope",
          "But: excited-state Hamiltonian and phonon rates not settled"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
c2 = card(s, MARGIN + 5.00, BODY_TOP, 4.72, 2.85, fill=TINT)
tf = c2.text_frame
run(tf.paragraphs[0], "NV — where we can compute", size=18, bold=True,
    color=NAVY)
for t in ["³E Hamiltonian measured and published",
          "Phonon-induced orbital hopping rate Γₓᵧ(T) measured",
          "Cryogenic ensemble EIT already demonstrated",
          "So the theory can be tested against something known"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
tb, tf = textbox(s, MARGIN, 4.60, BODY_W, 0.66)
p = tf.paragraphs[0]
rich(p, [("Strategy: ", {"bold": True, "color": NAVY}),
         ("build the criterion on NV, where every input is known, then carry "
          "the same criterion back to group-IV.", {})], size=16)
bottom(s, "NV wins as a test bed not on material grounds but because its "
          "excited-state model and phonon rates are independently measured.")
refs(s, ["[2] L. J. Rogers et al., Phys. Rev. Lett. 113, 263602 (2014).   "
         "[3] M. W. Doherty et al., Phys. Rep. 528, 1 (2013)."])

# ---- I-3 objective
s = add(L_CONTENT)
head(s, "We ask whether NV EIT survives to room temperature, and if not, "
        "exactly where it stops.")
title(s, "Introduction 3: The question")
q = card(s, MARGIN, BODY_TOP, BODY_W, 1.15, fill=NAVY)
tf = q.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run(p, "Is room-temperature NV EIT possible?", size=22, bold=True, color=WHITE)
p = para(tf)
p.alignment = PP_ALIGN.CENTER
run(p, "If not — up to what temperature can it be observed at all?",
    size=18, color=RGBColor(0xD8, 0xE2, 0xF3))
bullets(s, MARGIN, 2.95, BODY_W, 2.3, [
    [("Known: ", {"bold": True, "color": NAVY}),
     ("NV ensemble EIT is demonstrated at cryogenic temperature.", {})],
    [("Missing: ", {"bold": True, "color": NAVY}),
     ("any microscopic account of why it disappears on warming.", {})],
    [("The usual answer — “dephasing grows” — predicts only a monotone fade. "
      "It predicts no lower bound, no field threshold, and no change of "
      "sign.", {})],
    [("All three occur.", {"bold": True, "color": RED})],
])
bottom(s, "Cryogenic NV ensemble EIT is established; no published work fixes "
          "its upper temperature limit or its collapse mechanism.")
refs(s, ["[4] V. M. Acosta et al., Phys. Rev. Lett. 110, 213605 (2013).   "
         "[5] C. Santori et al., Phys. Rev. Lett. 97, 247401 (2006)."])

# ================================================== METHOD
section("Method", "II")

# ---- M-1 no-go / go
s = add(L_CONTENT)
head(s, "A no-go says the targeted pathway contributes nothing — not that the "
        "total response vanishes.")
title(s, "Method 1: What a no-go actually claims")
w1 = card(s, MARGIN, BODY_TOP, 4.72, 1.30, fill=BAD)
tf = w1.text_frame
run(tf.paragraphs[0], "✗   The tempting definition", size=17, bold=True,
    color=RED)
p = para(tf); p.space_before = Pt(5)
run(p, "“EIT is forbidden when χ_full = 0”", size=16, color=INK)
w2 = card(s, MARGIN + 5.00, BODY_TOP, 4.72, 1.30, fill=TINT)
tf = w2.text_frame
run(tf.paragraphs[0], "✓   The correct object", size=17, bold=True,
    color=NAVY)
p = para(tf); p.space_before = Pt(5)
run(p, "δχₛ ≡ χ_full − χ_cut  ≡  0", size=16, color=INK)
tb, tf = textbox(s, MARGIN, 3.04, BODY_W, 0.46)
rich(tf.paragraphs[0], [
    ("χ_full = 0 is what ", {}),
    ("perfect EIT", {"bold": True, "color": NAVY}),
    (" looks like.", {})], size=17)
c = card(s, MARGIN, 3.56, BODY_W, 1.66, fill=SOFT)
tf = c.text_frame
run(tf.paragraphs[0], "Ideal Λ system, worked through", size=17,
    bold=True, color=NAVY)
for t, v, em in [
        ("cut  (pathway removed)", "Ξ_cut = 1", False),
        ("full (pathway intact)", "Ξ_full = 0    ← response vanishes", False),
        ("sector correction", "δΞₛ = −1    ← maximally non-zero", True)]:
    p = para(tf); p.space_before = Pt(5)
    run(p, f"{t:<26s}", size=16, color=GREY, font="Courier New")
    run(p, v, size=16, bold=em, color=(NAVY if em else INK),
        font="Courier New")
bottom(s, "The pathway cut is an algebraic counterfactual (B = C = 0 with the "
          "steady state held fixed), not a control-off measurement.")
refs(s, ["[6] D. Finkelstein-Shapiro et al., Phys. Rev. A 99, 053829 (2019)."])

# ---- M-2 material independence
s = add(L_CONTENT)
head(s, "Built from the response kernel, the criterion holds for any "
        "Markovian emitter.")
title(s, "Method 2: A material-independent criterion")
steps = [("Material", "H, jump operators,\ndipoles"),
         ("Kernel", "Kᵢⱼ = dᵢ† G(z) dⱼ"),
         ("Test", "δχₛ = 0\n⇔  K₁₂K₂₁ = 0"),
         ("Verdict", "no-go / suppressed /\ngo")]
x = MARGIN
wid = (BODY_W - 3 * 0.32) / 4
for i, (t, d) in enumerate(steps):
    fill = NAVY if i == 3 else TINT
    fg = WHITE if i == 3 else NAVY
    c = card(s, x, BODY_TOP, wid, 1.55, fill=fill)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, t, size=18, bold=True, color=fg)
    for line in d.split("\n"):
        p = para(tf); p.alignment = PP_ALIGN.CENTER
        run(p, line, size=16, color=(WHITE if i == 3 else INK))
    if i < 3:
        cn = s.shapes.add_connector(1, Inches(x + wid + 0.05), Inches(2.40),
                                    Inches(x + wid + 0.27), Inches(2.40))
        cn.line.color.rgb = BLUE
        cn.line.width = Pt(2.25)
    x += wid + 0.32
bullets(s, MARGIN, 3.40, BODY_W, 1.85, [
    [("Nothing in the test is specific to NV, to diamond, or even to "
      "optics.", {})],
    [("Already applied to group-IV centres (SiV, SnV) and to a non-optical "
      "witness.", {})],
    [("What changes between materials is only which pathway is closed and by "
      "what.", {})],
])
bottom(s, "For the scalar block, δχₛ vanishes identically if and only if the "
          "product of the two Raman vertices K₁₂K₂₁ vanishes.")
refs(s, ["[7] K. Zhou et al., J. Phys. A 58, 095303 (2025).   "
         "[8] X. Zhao, L.-M. Kuang and J.-Q. Liao, "
         "Phys. Rev. A 113, 013723 (2026)."])

# ---- M-3 how we classify
s = add(L_CONTENT)
head(s, "We judge every point by three criteria at once, because a dip alone "
        "identifies nothing.")
title(s, "Method 3: How each operating point is judged")
crit = [("1. Sign", "Does the pathway reduce or\nincrease absorption?",
         "C > 0  transparency\nC < 0  induced absorption"),
        ("2. Containment", "Is the feature actually inside\nthe scan window?",
         "adaptive two-photon window\n(widths span 3 decades)"),
        ("3. Lineshape", "EIT interference, or\nAutler–Townes doublet?",
         "AIC model comparison\n|ΔAIC| ≥ 6")]
x = MARGIN
wid = (BODY_W - 2 * 0.30) / 3
for t, q_, a in crit:
    c = card(s, x, BODY_TOP, wid, 2.20, fill=TINT)
    tf = c.text_frame
    run(tf.paragraphs[0], t, size=18, bold=True, color=NAVY)
    for line in q_.split("\n"):
        p = para(tf); p.space_before = Pt(3)
        run(p, line, size=16, color=INK)
    for line in a.split("\n"):
        p = para(tf); p.space_before = Pt(3)
        run(p, line, size=16, color=GREY, italic=True)
    x += wid + 0.30
bullets(s, MARGIN, 4.00, BODY_W, 1.25, [
    [("Physics from the ", {}),
     ("full 9-level Lindblad generator", {"bold": True, "color": NAVY}),
     (" — no adiabatic elimination.", {})],
    [("240 grid points over T = 20–115 K and B⊥ = 0–0.5 T; each point gets "
      "one of four labels.", {})],
])
bottom(s, "A strong control gives an Autler–Townes doublet even with zero "
          "ground coherence, so a dip alone cannot establish EIT.")
refs(s, ["[9] P. M. Anisimov, J. P. Dowling and B. C. Sanders, "
         "Phys. Rev. Lett. 107, 163604 (2011)."])

# ================================================== RESULTS
section("Results", "III")

# ---- R-1 sanity check
s = add(L_CONTENT)
head(s, "The general formula reproduces the textbook three-level result as a "
        "one-line special case.")
title(s, "Result 1: Reduction to the textbook Λ system")
c = card(s, MARGIN, BODY_TOP, BODY_W, 1.70, fill=SOFT)
tf = c.text_frame
run(tf.paragraphs[0], "General (Schur complement, any number of excited "
                      "branches)", size=16, bold=True, color=NAVY)
p = para(tf); p.space_before = Pt(6)
run(p, "Ξ  =  S₁  −  β K₁₂K₂₁ / (γ_g + β S₂)", size=18, color=INK)
p = para(tf); p.space_before = Pt(8)
rich(p, [("One excited state, scalar dipoles ", {"color": GREY}),
         ("→", {"color": BLUE, "bold": True}),
         ("   Ξ = γ_g / (Aγ_g + β)   = the standard Λ-system result", {})],
     size=16)
bullets(s, MARGIN, 3.48, BODY_W, 1.80, [
    [("Recovers the published three-level susceptibility exactly — no "
      "adiabatic approximation is used anywhere.", {})],
    [("Holds for multi-branch excited manifolds: K₁₂ becomes a coherent sum "
      "over branches.", {})],
    [("Full Liouvillian reproduces the archived 70 K benchmark to ", {}),
     ("0.25 %", {"bold": True, "color": NAVY}), (".", {})],
])
bottom(s, "Textbook Λ susceptibility recovered as a special case; 70 K "
          "contrast C = 0.01387 against the archived 0.01384.")
refs(s, ["[1] M. Fleischhauer, A. Imamoğlu and J. P. Marangos, "
         "Rev. Mod. Phys. 77, 633 (2005).   Benchmark: [4]."])

# ---- R-2 phase diagram
s = add(L_CONTENT)
head(s, "NV EIT does not simply fade — it occupies a bounded island in the "
        "temperature–field plane.")
title(s, "Result 2: The temperature–field island")
picture(s, "sf_island.png", top=BODY_TOP, max_h=3.10)
# the cold edge is where this claim is weakest, so it is stated on the slide
# rather than left for the questions
tb, tf = textbox(s, MARGIN, 4.68, BODY_W, 0.60)
rich(tf.paragraphs[0], [
    ("Cold edge: ", {"bold": True, "color": NAVY}),
    ("30 K at 0.15 T; above 0.3 T it is the 20 K grid floor, closed by "
     "3 Autler–Townes points.", {})], size=16)
bottom(s, "Transparency needs B⊥ ≳ 0.15 T: the window runs 30–90 K at that "
          "field and 20–95 K above 0.3 T.")
refs(s, ["Phonon-induced orbital hopping: "
         "[10] K.-M. C. Fu et al., Phys. Rev. Lett. 103, 256404 (2009).   "
         "[11] M. L. Goldman et al., Phys. Rev. Lett. 114, 145502 (2015)."])

# ---- R-3 room temperature
s = add(L_CONTENT)
head(s, "Room temperature loses nine orders of contrast and lands two "
        "orders below the detection floor.")
title(s, "Result 3: Room temperature")
picture(s, "sf_contrast_T.png", top=BODY_TOP, max_h=3.66, max_w=6.80, cx=3.54)
rows = [("30 K", "0.99", "instant"), ("70 K", "1.4×10⁻²", "1 μs"),
        ("90 K", "1.5×10⁻⁴", "10 ms"),
        ("300 K", "−1.1×10⁻⁹", "never")]
c = card(s, 7.06, 1.70, 2.80, 3.40, fill=TINT)
tf = c.text_frame
run(tf.paragraphs[0], "contrast → time", size=17, bold=True, color=NAVY)
for t, cc, tau in rows:
    p = para(tf); p.space_before = Pt(7)
    col = RED if t == "300 K" else INK
    run(p, f"{t}   ", size=16, bold=True, color=col)
    run(p, f"{cc}", size=16, color=col)
    p2 = para(tf)
    run(p2, f"          {tau}", size=16, italic=True, color=GREY)
bottom(s, "At 300 K the contrast is 1.1×10⁻⁹, of the wrong sign, and 140× "
          "below the 1.5×10⁻⁷ detection floor.")
refs(s, ["³E structure: [3] M. W. Doherty et al., Phys. Rep. 528, 1 (2013).   "
         "Contrasts and detection chain: this work."])

# ---- R-4 sign reversal
s = add(L_CONTENT)
head(s, "Above 103 K the control stops inducing transparency and starts "
        "inducing absorption.")
title(s, "Result 4: Sign reversal at 103 K")
picture(s, "sf_lineshape.png", top=BODY_TOP, max_h=3.66, max_w=4.79, cx=2.54)
c = card(s, 5.14, BODY_TOP, 4.72, 1.72, fill=TINT)
tf = c.text_frame
run(tf.paragraphs[0], "Why this matters", size=17, bold=True, color=NAVY)
for t in ["No Autler–Townes mechanism gives a sign change",
          "So it directly tests the interference account"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
c = card(s, 5.14, 3.56, 4.72, 1.72, fill=SOFT)
tf = c.text_frame
run(tf.paragraphs[0], "Where it happens", size=17, bold=True, color=NAVY)
p = para(tf); p.space_before = Pt(6)
run(p, "T_sign = ", size=16, color=INK)
run(p, "103 [97, 108] K", size=18, bold=True, color=RED)
p = para(tf); p.space_before = Pt(4)
run(p, "full Liouvillian, 68 % interval, 80 samples", size=16, color=GREY)
p = para(tf); p.space_before = Pt(3)
run(p, "reduced kernel not trusted above 90 K", size=16, color=GREY)
bottom(s, "Peak contrast turns negative between 100 and 105 K above 0.15 T; "
          "at 105 K the transmission signal is −1.2×10⁻⁵.")
refs(s, ["Induced absorption in Λ systems: "
         "[12] P. Valente, H. Failache and A. Lezama, "
         "Phys. Rev. A 67, 013806 (2003)."])

# ---- R-5 transverse field
s = add(L_CONTENT)
head(s, "The transverse field is a switch that opens the pathway, not a dial "
        "that buys temperature.")
title(s, "Result 5: The role of the transverse field")
tbl = [("B⊥ (T)", "T₁% (K)", "verdict"),
       ("0 – 0.10", "27 – 33", "no usable window"),
       ("0.15", "70.3", "window opens"),
       ("0.2323", "71.1", "—"),
       ("0.50", "69.7", "no further gain")]
x0, y0 = MARGIN, BODY_TOP
cw = [1.70, 1.60, 2.55]
for r, row in enumerate(tbl):
    xx = x0
    for i, cell in enumerate(row):
        isf = (r == 0)
        c = card(s, xx, y0 + r * 0.68, cw[i], 0.62,
                 fill=(NAVY if isf else (TINT if r % 2 else WHITE)))
        if not isf and r % 2 == 0:
            c.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            c.line.width = Pt(0.75)
        tf = c.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run(p, cell, size=16, bold=isf or (i == 1 and not isf),
            color=(WHITE if isf else INK))
        xx += cw[i] + 0.10
bullets(s, 6.30, BODY_TOP + 0.10, 3.56, 3.4, [
    [("0 → 0.15 T", {"bold": True, "color": NAVY})],
    [("boundary jumps by 40 K", {})],
    [("0.15 → 0.50 T", {"bold": True, "color": NAVY})],
    [("boundary moves < 2 K", {"bold": True, "color": RED})],
], size=16, space=10)
bottom(s, "The field is required to have a channel at all; once open, the "
          "upper temperature is set by dissipation, not by the field.")
refs(s, ["Ground-state Zeeman structure: "
         "[3] M. W. Doherty et al., Phys. Rep. 528, 1 (2013).   "
         "Boundaries: this work."])

# ================================================== DISCUSSION
section("Discussion", "IV")

# ---- D-1 what to measure
s = add(L_CONTENT)
head(s, "The decisive measurement is following the fringe through zero at "
        "103 K.")
title(s, "Discussion 1: What to measure")
picture(s, "sf_observables.png", top=BODY_TOP, max_h=3.66, max_w=5.69,
        cx=2.99)
c1 = card(s, 6.02, BODY_TOP, 3.84, 1.72, fill=TINT)
tf = c1.text_frame
run(tf.paragraphs[0], "Where to look", size=17, bold=True, color=NAVY)
for t in ["30–90 K at 0.15 T, 20–95 K above 0.3 T",
          "Transmission or fluorescence both work"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
c2 = card(s, 6.02, 3.56, 3.84, 1.72, fill=WARN)
tf = c2.text_frame
run(tf.paragraphs[0], "The decisive measurement", size=17, bold=True,
    color=NAVY)
for t in ["Follow the fringe through zero at 103 K",
          "1.6 s at 105 K, 5.6 s at 110 K"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
bottom(s, "Transmission is detectable to 110 K, fluorescence to 105 K; past "
          "120 K technical noise blocks detection outright.")
refs(s, ["Ensemble EIT magnetometry, from which the detection chain is taken: "
         "[4] V. M. Acosta et al., Phys. Rev. Lett. 110, 213605 (2013)."])

# ---- D-2 why RT fails
s = add(L_CONTENT)
head(s, "Room temperature fails because a thermalising orbital mode destroys "
        "the Raman path.")
title(s, "Discussion 2: Why room temperature fails")
c = card(s, MARGIN, BODY_TOP, BODY_W, 1.55, fill=SOFT)
tf = c.text_frame
run(tf.paragraphs[0], "The mechanism", size=17, bold=True, color=NAVY)
p = para(tf); p.space_before = Pt(5)
rich(p, [("A jump inside the excited manifold conserves population but acts "
          "on every optical coherence as ", {}),
         ("pure loss", {"bold": True, "color": RED}),
         (".  Phonon hopping Eₓ ↔ Eᵧ therefore decoheres the pathway while "
          "removing no population.  Rate entering the coherence equation: "
          "Γₓᵧ / 4.", {})], size=16)
c1 = card(s, MARGIN, 3.32, 4.72, 1.96, fill=TINT)
tf = c1.text_frame
run(tf.paragraphs[0], "What NV would need", size=17, bold=True, color=NAVY)
for t in ["Freeze out the orbital dynamics, or",
          "Push the branch splitting far above kT",
          "Neither is available in bulk NV at 300 K"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
c2 = card(s, MARGIN + 5.00, 3.32, 4.72, 1.96, fill=SOFT)
tf = c2.text_frame
run(tf.paragraphs[0], "Back to group-IV", size=17, bold=True, color=GREY)
for t in ["Orbital Λ: no spin-overlap obstruction",
          "But orbital coherence damped by single phonons",
          "Strain and spin–orbit splitting are the levers"]:
    p = para(tf); p.space_before = Pt(6)
    run(p, "•  " + t, size=16, color=INK)
bottom(s, "The obstruction is a thermalising orbital degree of freedom in the "
          "excited manifold, not the spin dephasing usually blamed.")
refs(s, ["Orbital averaging and phonon coupling: "
         "[13] G. Thiering and A. Gali, Phys. Rev. B 96, 081115 (2017).   "
         "[14] K. D. Jahnke et al., New J. Phys. 17, 043011 (2015)."])

# ================================================== SUMMARY
s = add(L_CONTENT)
head(s, "We have fixed the operating window and found a falsifiable "
        "prediction to test this year.")
title(s, "Summary and next step")
findings = [
    ("Bounded island", "Transparency needs B⊥ ≳ 0.15 T, and runs 30–90 K "
                       "at that field, 20–95 K above 0.3 T."),
    ("Room temperature: no", "Residual contrast 1.1×10⁻⁹, wrong sign, "
                             "undetectable at any integration time."),
    ("Sign reversal at 103 K", "Control-induced absorption above "
                               "103 [97, 108] K — measurable in seconds."),
]
y = BODY_TOP
for i, (t, d) in enumerate(findings, 1):
    c = card(s, MARGIN, y, BODY_W, 0.92,
             fill=(TINT if i != 3 else WARN))
    tf = c.text_frame
    p = tf.paragraphs[0]
    run(p, f"{i}.  ", size=18, bold=True, color=NAVY)
    run(p, t, size=18, bold=True, color=NAVY)
    p2 = para(tf)
    run(p2, d, size=16, color=INK)
    y += 0.96
c = card(s, MARGIN, y, BODY_W, 0.78, fill=NAVY)
tf = c.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
run(p, "Next:  ", size=18, bold=True, color=RGBColor(0xD8, 0xE2, 0xF3))
run(p, "measure within the year and test the predicted window and the sign "
       "reversal against experiment.", size=18, bold=True, color=WHITE)
bottom(s, "PRA manuscript complete; six figures and all numerical gates "
          "recorded and reproducible from the repository.")
refs(s, ["All contrasts, boundaries and integration times: this work, "
         "full 9-level Lindblad generator."])

# ================================================== SUPPLEMENTARY
section("Appendix", "S")

# ---- A-1 theorems
s = add(L_CONTENT)
head(s, "The criterion rests on a small set of theorems, each provable with "
        "first-course linear algebra.")
title(s, "Appendix 1: Theorems behind the criterion")
th = [("1A", "Optical dark subspace has dimension N_g − rank Ω; basis "
             "independent."),
      ("1B", "Stationary iff it is an eigenvector of every jump operator and "
             "of H_eff."),
      ("2A", "The weak-probe response is an exact Schur complement."),
      ("2B", "δχₛ = 0 ⇔ K₁₂K₂₁ = 0.  This, not χ_full = 0, is the no-go."),
      ("3",  "Under fast damping the first non-zero path moment fixes the "
             "power of 1/Γ."),
      ("4",  "Finitely many moments certify an all-orders exact zero "
             "(Cayley–Hamilton)."),
      ("5",  "Sector symmetry forces an identically vanishing transfer."),
      ("6",  "Singular damping leaves an O(1) protected channel in ker D.")]
y = BODY_TOP
for k, d in th:
    tb, tf = textbox(s, MARGIN, y, BODY_W, 0.40)
    p = tf.paragraphs[0]
    run(p, f"Thm {k:<3s} ", size=16, bold=True, color=NAVY)
    run(p, d, size=16, color=INK)
    y += 0.45
bottom(s, "Full statements and proofs are in the group's theory note; the "
          "manuscript uses 2A, 2B and the rate map only.")
refs(s, ["Dark-state classification: "
         "[6] D. Finkelstein-Shapiro et al., Phys. Rev. A 99, 053829 (2019).  "
         " Morris–Shore: [15] J. R. Morris and B. W. Shore, "
         "Phys. Rev. A 27, 906 (1983)."])

# ---- A-2 model audit
s = add(L_CONTENT)
head(s, "The reduced kernel reproduces boundary temperatures but not contrast "
        "magnitudes.")
title(s, "Appendix 2: Reduced kernel vs full Liouvillian")
bullets(s, MARGIN, BODY_TOP, BODY_W, 1.80, [
    [("Compared against the full Liouvillian on 108 grid points: ", {}),
     ("47 agree within 10 %", {"bold": True}),
     (", and ", {}), ("13 disagree in sign", {"bold": True, "color": RED}),
     (".", {})],
    [("Failures are structured, not random: below 40 K the reduced kernel "
      "reports transparency where the full model reports absorption; above "
      "90 K no field agrees within 10 %.", {})],
    [("Threshold temperatures still agree — the bands overlap, medians "
      "shifted by ≈ 3 K.", {})],
])
c = card(s, MARGIN, 3.50, BODY_W, 1.78, fill=TINT)
tf = c.text_frame
run(tf.paragraphs[0], "How the manuscript uses this", size=17, bold=True,
    color=NAVY)
for t in ["Every contrast magnitude quoted comes from the full Liouvillian",
          "The reduced kernel is used only for threshold locations, "
          "re-derived with the full model as a check",
          "Truncation (singlet, hyperfine) shifts contrast by ≤ 30 %, never "
          "the sign"]:
    p = para(tf); p.space_before = Pt(5)
    run(p, "•  " + t, size=16, color=INK)
bottom(s, "Reduced-vs-full agreement is 47/108 points within 10 % with 13 "
          "sign disagreements, concentrated below 40 K and above 90 K.")
refs(s, ["Full 9-level generator built from "
         "[3] M. W. Doherty et al., Phys. Rep. 528, 1 (2013)."])

# ---- A-3 B-perp exponent
s = add(L_CONTENT)
head(s, "The quadratic law is confirmed at the warm end of the island, not at "
        "the cold end.")
title(s, "Appendix 3: Transverse-field exponent")
picture(s, "sf_bperp.png", max_h=BODY_BOT - BODY_TOP)
bottom(s, "85 K gives n = 2.11 ± 0.08, stable against the fitting cutoff; at "
          "55 and 70 K the exponent drifts and is not determined.")
refs(s, ["Expected C = C_res + a B⊥² from one symmetry-breaking insertion per "
         "Λ leg; ground-state mixing scale γₑB = D_gs from [3]."])

# ---- A-4 EIT vs ATS
s = add(L_CONTENT)
head(s, "We separate EIT from Autler–Townes by model comparison rather than "
        "by eye.")
title(s, "Appendix 4: EIT vs Autler–Townes")
c = card(s, MARGIN, BODY_TOP, 4.72, 2.05, fill=TINT)
tf = c.text_frame
run(tf.paragraphs[0], "EIT model (interference)", size=17, bold=True,
    color=NAVY)
p = para(tf); p.space_before = Pt(7)
run(p, "A = C₊²/(γ₊²+δ²) − C₋²/(γ₋²+δ²)", size=16, color=INK)
p = para(tf); p.space_before = Pt(7)
run(p, "a narrow window inside a broad line", size=16, italic=True, color=GREY)
c = card(s, MARGIN + 5.00, BODY_TOP, 4.72, 2.05, fill=SOFT)
tf = c.text_frame
run(tf.paragraphs[0], "ATS model (doublet)", size=17, bold=True, color=GREY)
p = para(tf); p.space_before = Pt(7)
run(p, "A = C²[1/(γ²+(δ−δ₀)²) + 1/(γ²+(δ+δ₀)²)]", size=16, color=INK)
p = para(tf); p.space_before = Pt(7)
run(p, "two dressed-state peaks, no interference", size=16, italic=True,
    color=GREY)
bullets(s, MARGIN, 3.85, BODY_W, 1.43, [
    [("Both fitted to the same computed spectrum, together with Lorentzian "
      "and Fano alternatives.", {})],
    [("Verdict from ΔAIC = IC_ATS − IC_EIT with a fixed robust gate at "
      "|ΔAIC| ≥ 6, unchanged under 60 noise bootstraps and 30 randomised "
      "fit initialisations.", {})],
])
bottom(s, "Following Anisimov, Dowling and Sanders, PRL 107, 163604 (2011); "
          "the classification is stable at every temperature tested.")
refs(s, ["[9] P. M. Anisimov, J. P. Dowling and B. C. Sanders, "
         "Phys. Rev. Lett. 107, 163604 (2011)."])

# ---- A-5 detection chain
s = add(L_CONTENT)
head(s, "One detection chain at every temperature makes the comparison "
        "measure physics, not geometry.")
title(s, "Appendix 5: Detection chain")
bullets(s, MARGIN, BODY_TOP, BODY_W, 1.50, [
    [("The sector cross-section grows as the optical line narrows: a fixed "
      "sample is OD ≈ 16 at 10 K and OD ≈ 0.04 at 300 K.", {})],
    [("Primary axis is therefore the ", {}),
     ("OD-matched sample", {"bold": True, "color": NAVY}),
     (" (sector OD = 1); the fixed 1 ppm / 0.5 mm sample is reported "
      "alongside.", {})],
])
par = [("λ", "637 nm"), ("Debye–Waller", "0.030"),
       ("γ_inh", "30 GHz"), ("probe power", "1 μW"),
       ("η detect", "0.1"), ("σ_tech", "10⁻⁶"),
       ("target SNR", "5"), ("ceiling", "24 h")]
x, y = MARGIN, 3.30
for i, (k, v) in enumerate(par):
    c = card(s, x, y, 2.34, 0.86, fill=TINT)
    tf = c.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    run(p, k + "   ", size=16, color=GREY)
    run(p, v, size=16, bold=True, color=NAVY)
    x += 2.46
    if (i + 1) % 4 == 0:
        x = MARGIN
        y += 0.98
bottom(s, "Where the archived conversions disagreed (70 K and 300 K), the "
          "conservative value was taken and the disagreement recorded.")
refs(s, ["Debye–Waller factor, inhomogeneous width and ensemble geometry: "
         "[4] V. M. Acosta et al., Phys. Rev. Lett. 110, 213605 (2013)."])

prs.save(str(OUT))
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

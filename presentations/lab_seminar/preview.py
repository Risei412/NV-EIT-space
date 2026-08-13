"""Approximate raster preview of the deck.

LibreOffice will not run here, so this draws the slide geometry directly with
PIL: filled shapes, embedded pictures, and word-wrapped text at the declared
font sizes.  It is not a faithful renderer -- font metrics and autofit differ
from PowerPoint -- but it shows composition, alignment and gross overflow,
which is what the visual pass is for.
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

EMU = 914400.0
DPI = 96
SW, SH = 10.0, 7.5

FONT_DIRS = ["/usr/share/fonts/truetype/dejavu",
             "/usr/share/fonts/truetype/liberation"]


def load_font(size_px, bold=False, italic=False):
    names = (["LiberationSans-BoldItalic.ttf", "DejaVuSans-BoldOblique.ttf"]
             if bold and italic else
             ["LiberationSans-Bold.ttf", "DejaVuSans-Bold.ttf"] if bold else
             ["LiberationSans-Italic.ttf", "DejaVuSans-Oblique.ttf"]
             if italic else
             ["LiberationSans-Regular.ttf", "DejaVuSans.ttf"])
    for d in FONT_DIRS:
        for n in names:
            p = Path(d) / n
            if p.exists():
                try:
                    return ImageFont.truetype(str(p), size_px)
                except Exception:
                    pass
    return ImageFont.load_default()


def rgb(c, default=(30, 30, 30)):
    try:
        if c and c.type is not None and c.rgb is not None:
            return tuple(c.rgb)
    except Exception:
        pass
    return default


def wrap(draw, text, font, max_w):
    out, line = [], ""
    for word in text.split():
        t = (line + " " + word).strip()
        if draw.textlength(t, font=font) <= max_w or not line:
            line = t
        else:
            out.append(line)
            line = word
    if line:
        out.append(line)
    return out


def draw_shape(img, draw, sh):
    if sh.left is None:
        return
    x, y = sh.left / EMU * DPI, sh.top / EMU * DPI
    w, h = (sh.width or 0) / EMU * DPI, (sh.height or 0) / EMU * DPI

    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
            im = im.resize((max(1, int(w)), max(1, int(h))))
            img.paste(im, (int(x), int(y)), im)
        except Exception:
            draw.rectangle([x, y, x + w, y + h], outline=(150, 150, 150))
        return

    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        fill = None
        try:
            if sh.fill.type is not None and str(sh.fill.type) != "BACKGROUND (5)":
                fill = rgb(sh.fill.fore_color, (245, 245, 245))
        except Exception:
            pass
        if fill:
            draw.rounded_rectangle([x, y, x + w, y + h], radius=6, fill=fill)

    if not sh.has_text_frame:
        return
    tf = sh.text_frame
    ml = (tf.margin_left.inches if tf.margin_left is not None else 0.1) * DPI
    mr = (tf.margin_right.inches if tf.margin_right is not None else 0.1) * DPI
    mt = (tf.margin_top.inches if tf.margin_top is not None else 0.05) * DPI
    cx = x + ml
    cy = y + mt
    avail = max(10, w - ml - mr)

    lines = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            continue
        size = max((r.font.size.pt if r.font.size else 16) for r in runs)
        bold = any(r.font.bold for r in runs)
        ital = any(r.font.italic for r in runs)
        col = rgb(runs[0].font.color)
        f = load_font(int(size / 72 * DPI), bold, ital)
        sb = (p.space_before.pt if p.space_before else 0) / 72 * DPI
        sa = (p.space_after.pt if p.space_after else 0) / 72 * DPI
        txt = "".join(r.text for r in runs)
        for i, ln in enumerate(wrap(draw, txt, f, avail)):
            lines.append((ln, f, col, size / 72 * DPI * 1.22,
                          sb if i == 0 else 0, 0))
        if lines:
            ln = list(lines[-1]); ln[5] = sa; lines[-1] = tuple(ln)

    total = sum(lh + sb + sa for _, _, _, lh, sb, sa in lines)
    anchor = str(tf.vertical_anchor)
    if "MIDDLE" in anchor:
        cy = y + (h - total) / 2
    for ln, f, col, lh, sb, sa in lines:
        cy += sb
        draw.text((cx, cy), ln, font=f, fill=col)
        cy += lh + sa


def main(path, out_dir, only=None):
    prs = Presentation(path)
    out = Path(out_dir)
    out.mkdir(exist_ok=True)
    made = []
    for n, slide in enumerate(prs.slides, 1):
        if only and n not in only:
            continue
        img = Image.new("RGB", (int(SW * DPI), int(SH * DPI)), "white")
        draw = ImageDraw.Draw(img)
        for sh in slide.slide_layout.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                draw_shape(img, draw, sh)
        for sh in slide.shapes:
            draw_shape(img, draw, sh)
        draw.rectangle([0, 0, int(SW * DPI) - 1, int(SH * DPI) - 1],
                       outline=(200, 200, 200))
        f = out / f"s{n:02d}.png"
        img.save(f)
        made.append(str(f))
    print("\n".join(made))


if __name__ == "__main__":
    only = {int(a) for a in sys.argv[3:]} if len(sys.argv) > 3 else None
    main(sys.argv[1], sys.argv[2], only)

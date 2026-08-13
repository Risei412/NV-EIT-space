"""Geometric QA for the deck.

LibreOffice cannot render in this environment, so the usual look-at-the-pixels
pass is unavailable.  This substitutes the checks that catch the defects that
pass would have caught: shapes off the canvas or too close to its edge,
overlapping shapes, and text that cannot fit the box it is in.

Line counts come from real Liberation Sans metrics, which are compatible with
the Arial metrics PowerPoint will use, at 1.22 line spacing.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU = 914400.0
SW, SH = 10.0, 7.5
EDGE_MIN = 0.30          # template art legitimately sits closer than 0.45
OVERLAP_TOL = 0.02       # inches of tolerated overlap

# Line counts come from real Liberation Sans metrics (metric-compatible with
# Arial) rather than an average character width -- the flat average called a
# third line on almost every two-line message.
LINE_H = 1.22
_REF_PT = 200
_FDIR = "/usr/share/fonts/truetype/liberation"
_FACES = {(False, False): "LiberationSans-Regular.ttf",
          (True, False): "LiberationSans-Bold.ttf",
          (False, True): "LiberationSans-Italic.ttf",
          (True, True): "LiberationSans-BoldItalic.ttf"}
_CACHE = {}


def _font(bold, italic):
    key = (bold, italic)
    if key not in _CACHE:
        from PIL import ImageFont
        _CACHE[key] = ImageFont.truetype(f"{_FDIR}/{_FACES[key]}", _REF_PT)
    return _CACHE[key]


def _draw():
    if "d" not in _CACHE:
        from PIL import Image, ImageDraw
        _CACHE["d"] = ImageDraw.Draw(Image.new("RGB", (8, 8)))
    return _CACHE["d"]


def wrapped_lines(text, size_pt, width_in, bold=False, italic=False):
    """Number of rendered lines for word-wrapped text in a box of width_in."""
    if not text.strip():
        return 1
    d, f = _draw(), _font(bold, italic)
    limit = width_in * 72.0 * _REF_PT / size_pt
    n, cur = 1, ""
    for word in text.split():
        cand = (cur + " " + word).strip()
        if d.textlength(cand, font=f) <= limit or not cur:
            cur = cand
        else:
            n += 1
            cur = word
    return n


# template decorations we did not place and should not police
TEMPLATE_NAMES = ("正方形/長方形", "図 ", "スライド番号", "Picture Placeholder")


def is_template_art(sh):
    return any(sh.name.startswith(t) or t in sh.name for t in TEMPLATE_NAMES)


def box(sh):
    return (sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU)


def text_height(tf, width_in):
    h = 0.0
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        size = max((r.font.size.pt if r.font.size else 16) for r in runs) \
            if runs else 16
        text = "".join(r.text for r in runs)
        bold = any(r.font.bold for r in runs)
        ital = any(r.font.italic for r in runs)
        n = wrapped_lines(text, size, width_in, bold, ital) if text else 1
        h += n * LINE_H * size / 72.0
        if p.space_before is not None:
            h += p.space_before.pt / 72.0
        if p.space_after is not None:
            h += p.space_after.pt / 72.0
    return h


def rects_overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = min(ax + aw, bx + bw) - max(ax, bx)
    oy = min(ay + ah, by + bh) - max(ay, by)
    return ox > OVERLAP_TOL and oy > OVERLAP_TOL, ox, oy


def main(path):
    prs = Presentation(path)
    issues = 0
    for n, slide in enumerate(prs.slides, 1):
        mine = []
        for sh in slide.shapes:
            if is_template_art(sh):
                continue
            if sh.left is None or sh.top is None:
                continue
            b = box(sh)
            mine.append((sh, b))

            # --- off-canvas / edge margin
            x, y, w, h = b
            if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
                print(f"[{n:2d}] OFF-CANVAS  {sh.shape_type} '{sh.name[:22]}' "
                      f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
                issues += 1
            elif x < EDGE_MIN or y < EDGE_MIN or (SW - (x + w)) < EDGE_MIN \
                    or (SH - (y + h)) < EDGE_MIN:
                print(f"[{n:2d}] TIGHT-EDGE  '{sh.name[:22]}' "
                      f"x={x:.2f} y={y:.2f} r={SW-(x+w):.2f} b={SH-(y+h):.2f}")
                issues += 1

            # --- text fit
            if sh.has_text_frame and sh.text_frame.text.strip():
                tf = sh.text_frame
                ml = tf.margin_left.inches if tf.margin_left else 0.1
                mr = tf.margin_right.inches if tf.margin_right else 0.1
                mt = tf.margin_top.inches if tf.margin_top else 0.05
                mb = tf.margin_bottom.inches if tf.margin_bottom else 0.05
                avail_w = w - ml - mr
                avail_h = h - mt - mb
                need = text_height(tf, max(avail_w, 0.2))
                if need > avail_h * 1.06:
                    print(f"[{n:2d}] TEXT-OVERFLOW '{sh.name[:20]}' "
                          f"needs {need:.2f}in in {avail_h:.2f}in  "
                          f"({sh.text_frame.text[:52]!r})")
                    issues += 1
                # font-size floor
                for p in tf.paragraphs:
                    for r in p.runs:
                        if r.font.size and r.font.size.pt < 16 and r.text.strip():
                            print(f"[{n:2d}] SMALL-FONT {r.font.size.pt:.0f}pt "
                                  f"'{r.text[:34]}'")
                            issues += 1

        # --- overlaps among the shapes we placed
        for i in range(len(mine)):
            for j in range(i + 1, len(mine)):
                (s1, b1), (s2, b2) = mine[i], mine[j]
                # a picture inside a card is fine; connectors are lines
                if "Connector" in s1.name or "Connector" in s2.name:
                    continue
                ov, ox, oy = rects_overlap(b1, b2)
                if ov:
                    print(f"[{n:2d}] OVERLAP '{s1.name[:18]}' x "
                          f"'{s2.name[:18]}'  {ox:.2f} x {oy:.2f} in")
                    issues += 1
    print(f"\n{issues} geometric issue(s) across {len(prs.slides)} slides")
    return issues


if __name__ == "__main__":
    sys.exit(0 if main(sys.argv[1]) == 0 else 1)
